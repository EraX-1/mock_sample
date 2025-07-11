import logging
import os
import datetime
import uuid
import subprocess
import tempfile
import io
import base64
import re
import traceback
from concurrent.futures import ThreadPoolExecutor, as_completed
from tqdm import tqdm

import azure.functions as azfunc
from azure.storage.blob import BlobServiceClient

# 外部ライブラリ
import html2text
from bs4 import BeautifulSoup
from docx import Document
from openpyxl import load_workbook, Workbook
from pptx import Presentation
from pypdf import PdfReader, PdfWriter
import img2pdf
from PIL import Image

# ====================================================
# サービス設定（ハードコード版）
# ====================================================
# Azure AI Document Intelligence
from azure.core.credentials import AzureKeyCredential
from azure.ai.documentintelligence import DocumentIntelligenceClient
from langchain_community.document_loaders import AzureAIDocumentIntelligenceLoader

# Azure SDKのログレベルをWARNINGに設定
logging.getLogger("azure").setLevel(logging.WARNING)

from openai import AzureOpenAI as AOAI


class AzureOpenAI:
    def __init__(self):
        self.aoai_api_key = os.getenv("AZURE_OPENAI_API_KEY")
        self.aoai_endpoint = os.getenv(
            "AZURE_OPENAI_ENDPOINT",
            "https://aoai-kdk-knowledgedb-dev-eastus-001.openai.azure.com/"
        )
        self.aoai_api_version = os.getenv("AZURE_OPENAI_API_VERSION", "2024-09-01-preview")

    def init_client(self):
        client = AOAI(
            api_key=self.aoai_api_key,
            api_version=self.aoai_api_version,
            azure_endpoint=self.aoai_endpoint,
        )
        return client


class AzureAIDocumentIntelligence:
    def __init__(self):
        self.document_intelligence_key = "7sBJx2laLkeunTVQwo5W9K545V8xO1DivZTXBk0yIxO8ocyUuHwEJQQJ99BAAC1i4TkXJ3w3AAALACOG0eBn"
        self.document_intelligence_endpoint = (
            "https://di-kdk-knowledgedb-dev-cus-001.cognitiveservices.azure.com/"
        )

    def init_form_recognizer_client(self):
        return DocumentIntelligenceClient(
            endpoint=self.document_intelligence_endpoint,
            credential=AzureKeyCredential(self.document_intelligence_key),
        )

    def init_loader(self, file_path: str):
        return AzureAIDocumentIntelligenceLoader(
            file_path=file_path,
            api_key=self.document_intelligence_key,
            api_endpoint=self.document_intelligence_endpoint,
            api_model="prebuilt-layout",
            mode="markdown",
        )


# Azure AI Search
from azure.search.documents import SearchClient
from azure.search.documents.indexes import SearchIndexClient


class AzureAISearch:
    def __init__(self):
        self.search_endpoint = os.getenv(
            "AZURE_SEARCH_ENDPOINT",
            "https://srch-kdk-knoledgedb-dev-cus-001.search.windows.net"
        )
        self.search_api_key = os.getenv("AZURE_SEARCH_API_KEY")

    def init_search_index_client(self):
        search_index_client = SearchIndexClient(
            endpoint=self.search_endpoint,
            credential=AzureKeyCredential(self.search_api_key),
        )
        return search_index_client

    def init_search_client(self, index_name: str):
        search_client = SearchClient(
            endpoint=self.search_endpoint,
            index_name=index_name,
            credential=AzureKeyCredential(self.search_api_key),
        )
        return search_client

    def get_index_name(self, index_type: str):
        # インデックス名は、インデックスタイプ（フォルダ名）そのものとする
        return index_type


# Azure Blob Storage
class AzureBlobStorage:
    def __init__(self):
        self.connection_string = os.getenv("AZURE_STORAGE_CONNECTION_STRING")
        self.container_name = "auto-indexed-files"
        self.blob_service_endpoint = (
            "https://stkdkknowledgedbdocs001.blob.core.windows.net"
        )
        self.blob_service_client = BlobServiceClient.from_connection_string(
            self.connection_string
        )

    def init_container_client(self):
        container_client = self.blob_service_client.get_container_client(
            self.container_name
        )
        return container_client

    def get_blob_url(self, blob_name: str):
        return f"{self.blob_service_endpoint}/{self.container_name}/{blob_name}"


# ====================================================
# 環境設定（ハードコード版）
# ====================================================
DATABASE_URL = os.getenv("DATABASE_URL")
AZURE_STORAGE_CONNECTION_STRING = os.getenv("AZURE_STORAGE_CONNECTION_STRING")
EMBEDDING_MODEL_NAME = "text-embedding-ada-002"
SOURCE_CONTAINER = "sample-indexing-target-files"
DESTINATION_CONTAINER = "auto-indexed-files"

# ====================================================
# Database Setup (SQLAlchemy) - モデル定義
# ====================================================
from sqlalchemy import (
    create_engine,
    Column,
    Integer,
    String,
    DateTime,
    func as sql_func,
)
from sqlalchemy.ext.declarative import declarative_base
from sqlalchemy.orm import sessionmaker

engine = create_engine(DATABASE_URL)
SessionLocal = sessionmaker(autocommit=False, autoflush=False, bind=engine)
Base = declarative_base()


class IndexedFile(Base):
    __tablename__ = "indexed_files"
    id = Column(Integer, primary_key=True, index=True)
    original_blob_name = Column(String, unique=True, nullable=False, index=True)
    indexed_blob_name = Column(String, nullable=False)
    file_type = Column(String, nullable=False)
    # index_type は SearchIndexType.id を参照する
    index_type = Column(String(36), nullable=False)
    indexed_at = Column(DateTime, default=datetime.datetime.utcnow)


class SearchIndexType(Base):
    __tablename__ = "search_index_types"
    id = Column(String(36), primary_key=True, default=lambda: str(uuid.uuid4()))
    folder_name = Column(String(255), unique=True, nullable=False)
    created_at = Column(
        DateTime, server_default=sql_func.current_timestamp(), nullable=False
    )


Base.metadata.create_all(bind=engine)


# ====================================================
# Conversion Functions
# ====================================================
def convert_doc_bytes_to_docx_bytes(doc_bytes: bytes) -> bytes:
    with tempfile.TemporaryDirectory() as tmpdir:
        tmp_doc_path = os.path.join(tmpdir, "temp_input.doc")
        tmp_docx_path = os.path.join(tmpdir, "temp_input.docx")
        with open(tmp_doc_path, "wb") as f:
            f.write(doc_bytes)
        command = [
            "libreoffice",
            "--headless",
            "--convert-to",
            "docx",
            "--outdir",
            tmpdir,
            tmp_doc_path,
        ]
        try:
            subprocess.run(command, check=True)
            with open(tmp_docx_path, "rb") as f:
                docx_bytes = f.read()
            return docx_bytes
        except subprocess.CalledProcessError as e:
            raise Exception(f"変換中にエラーが発生しました: {str(e)}")


def convert_image_to_pdf(file: bytes) -> bytes:
    img = Image.open(io.BytesIO(file))
    fmt = img.format.lower()
    if fmt not in ["png", "jpeg", "jpg"]:
        raise ValueError(f"サポートされていない画像形式です: {fmt}")
    try:
        pdf_bytes = img2pdf.convert(file)
    except Exception as e:
        raise Exception(f"PDF生成中にエラーが発生しました: {str(e)}")
    return pdf_bytes


# ====================================================
# Extraction Functions
# ====================================================
def extract_markdown_text_from_docx(file: bytes) -> list[dict[str, any]]:
    paragraph_contents = []

    def process_paragraph(i_para):
        i, paragraph = i_para
        try:
            text = paragraph.text.strip()
            if text:
                markdown_text = ""
                if paragraph.style.name.startswith("Heading"):
                    level = int(paragraph.style.name[-1])
                    markdown_text = "#" * level + " " + text + "\n\n"
                else:
                    markdown_text = text + "\n\n"
                return {"content": markdown_text, "page": 1, "order": i}
            return None
        except Exception as e:
            raise Exception(f"段落の処理中にエラーが発生しました {i}: {e}")

    byte_stream = io.BytesIO(file)
    doc = Document(byte_stream)
    for i, para in enumerate(doc.paragraphs):
        res = process_paragraph((i, para))
        if res is not None:
            paragraph_contents.append(res)
    paragraph_contents.sort(key=lambda x: x["order"])
    from itertools import groupby

    docs = []
    for page, group in groupby(paragraph_contents, key=lambda x: x["page"]):
        combined_content = "".join(item["content"] for item in group)
        docs.append(
            {
                "page_content": combined_content + f"<PAGE_NUMBER>{page}</PAGE_NUMBER>",
                "metadata": {"page": page},
            }
        )
    return docs


def extract_markdown_text_from_pptx(file: bytes) -> list[dict[str, any]]:
    docs = []

    def process_slide(i_slide):
        i, slide = i_slide
        slide_docs = []
        try:
            texts = []
            if slide.shapes.title:
                texts.append(f"# {slide.shapes.title.text}\n")
            for shape in slide.shapes:
                if (
                    hasattr(shape, "text")
                    and shape.text.strip()
                    and shape != slide.shapes.title
                ):
                    text = shape.text.strip()
                    if text:
                        texts.append(text + "\n")
            content = "\n".join(texts).strip()
            if content:
                slide_docs.append(
                    {"page_content": content, "metadata": {"page": i + 1}}
                )
            return slide_docs
        except Exception as e:
            raise Exception(f"スライドの処理中にエラーが発生しました {i}: {e}")

    temp_file = tempfile.NamedTemporaryFile(delete=False)
    try:
        temp_file.write(file)
        temp_file.close()
        prs = Presentation(temp_file.name)
        for i, slide in enumerate(prs.slides):
            docs.extend(process_slide((i, slide)))
    finally:
        os.unlink(temp_file.name)
    docs = sorted(docs, key=lambda x: x["metadata"]["page"])
    for doc in docs:
        doc["page_content"] += (
            "<PAGE_NUMBER>" + str(doc["metadata"]["page"]) + "</PAGE_NUMBER>"
        )
    return docs


def extract_markdown_text_from_html(file: bytes) -> list[dict[str, any]]:
    docs = []

    def process_html_section(i_section):
        i, section = i_section
        section_docs = []
        try:
            h2t = html2text.HTML2Text()
            h2t.ignore_links = False
            h2t.ignore_images = False
            h2t.ignore_tables = False
            markdown_content = h2t.handle(str(section))
            if markdown_content.strip():
                section_docs.append(
                    {"page_content": markdown_content, "metadata": {"page": i + 1}}
                )
            return section_docs
        except Exception as e:
            raise Exception(f"セクションの処理中にエラーが発生しました {i}: {e}")

    try:
        soup = BeautifulSoup(file.decode("utf-8"), "html.parser")
        sections = soup.find_all(
            ["article", "section", "div"], class_=True, recursive=False
        )
        if not sections:
            sections = [soup.body] if soup.body else [soup]
        for i, section in enumerate(sections):
            docs.extend(process_html_section((i, section)))
    except Exception as e:
        raise Exception(f"Error processing HTML file: {e}")
    docs = sorted(docs, key=lambda x: x["metadata"]["page"])
    for doc in docs:
        doc["page_content"] += f"<PAGE_NUMBER>{doc['metadata']['page']}</PAGE_NUMBER>"
    return docs


def extract_markdown_text_from_pdf(file: bytes) -> list[dict[str, any]]:
    docs = []

    def process_page(i_page):
        i, page = i_page
        page_docs = []
        temp_page_file = None
        try:
            temp_page_file = tempfile.NamedTemporaryFile(delete=False)
            writer = PdfWriter()
            writer.add_page(page)
            bytes_io = io.BytesIO()
            writer.write(bytes_io)
            page_bytes = bytes_io.getvalue()
            bytes_io.close()
            temp_page_file.write(page_bytes)
            temp_page_file.close()
            loader = AzureAIDocumentIntelligence().init_loader(
                file_path=temp_page_file.name
            )
            doc = loader.load()
            for d in doc:
                page_docs.append(
                    {
                        "page_content": d.page_content,
                        "metadata": {"page": i + 1},
                    }
                )
            return page_docs
        except Exception as e:
            raise Exception(f"Error processing page {i}: {e}")
        finally:
            if temp_page_file is not None:
                os.unlink(temp_page_file.name)

    temp_file = tempfile.NamedTemporaryFile(delete=False)
    try:
        temp_file.write(file)
        temp_file.close()
        reader = PdfReader(temp_file.name)
        for i, page in enumerate(reader.pages):
            docs.extend(process_page((i, page)))
    finally:
        os.unlink(temp_file.name)
    docs = sorted(docs, key=lambda x: x["metadata"]["page"])
    for doc in docs:
        doc["page_content"] += (
            "<PAGE_NUMBER>" + str(doc["metadata"]["page"]) + "</PAGE_NUMBER>"
        )
    return docs


def extract_markdown_text_from_image(file: bytes) -> list[dict[str, any]]:
    pdf_bytes = convert_image_to_pdf(file)
    docs = extract_markdown_text_from_pdf(pdf_bytes)
    return docs


def extract_markdown_text_from_excel(file: bytes):
    excel_file = io.BytesIO(file)
    wb = load_workbook(excel_file, keep_vba=False)
    wb.vba_archive = None
    wb.is_template = False
    wb.vba_modified = False
    docs = []
    for i, sheet in enumerate(wb.worksheets):
        logging.info(f"シート名: {sheet.title}")
        new_wb = Workbook()
        new_ws = new_wb.active
        new_ws.title = sheet.title
        none_count = 0
        max_rows = sheet.max_row
        max_cols = sheet.max_column
        for row_idx in tqdm(
            range(1, max_rows + 1), desc=f"Processing sheet {sheet.title}"
        ):
            none_row = True
            for col_idx in range(1, max_cols + 1):
                cell = sheet.cell(row=row_idx, column=col_idx)
                if cell.value is not None:
                    none_row = False
                new_ws.cell(row=row_idx, column=col_idx, value=cell.value)
            if none_row:
                none_count += 1
            else:
                none_count = 0
            if none_count > 10:
                break
        sheet_buffer = io.BytesIO()
        new_wb.save(sheet_buffer)
        sheet_buffer.seek(0)
        with tempfile.NamedTemporaryFile(delete=False, suffix=".xlsx") as temp_file:
            temp_file.write(sheet_buffer.getvalue())
            temp_file.flush()
            loader = AzureAIDocumentIntelligence().init_loader(file_path=temp_file.name)
            doc = loader.load()
            docs.append(
                {
                    "page_content": doc[0].page_content,
                    "metadata": {"page": i + 1},
                }
            )
    return docs


# ====================================================
# Indexing Functions
# ====================================================
def _encode_data(data: str) -> str:
    return base64.urlsafe_b64encode(data.encode("utf-8")).decode("utf-8")


def _index_docs_to_azure_ai_search(
    chunks,
    source_file_name: str,
    index_type_id: str,
    indexed_blob_name: str,
    batch_size=5,
):
    # index_type_idからfolder_nameを取得してインデックス名として使用
    db = SessionLocal()
    try:
        index_type = (
            db.query(SearchIndexType)
            .filter(SearchIndexType.id == index_type_id)
            .first()
        )
        if not index_type:
            raise Exception(f"No search index type found for id: {index_type_id}")
        # UUIDをそのままインデックス名として使用
        index_name = index_type_id
    finally:
        db.close()

    search_client = AzureAISearch().init_search_client(index_name)
    open_ai_client = AzureOpenAI().init_client()

    def _embed_and_prepare_document(
        open_ai_client, chunk, source_file_name, indexed_blob_name, i
    ):
        try:
            response = open_ai_client.embeddings.create(
                input=chunk["content"], model=EMBEDDING_MODEL_NAME
            )
            document = {
                "id": _encode_data(source_file_name + "_" + str(i)),
                "keywords": chunk["keywords"],
                "content": chunk["content"],
                "contentVector": response.data[0].embedding,
                "pageNumber": chunk["page_number"],
                "sourceFileName": source_file_name,
                "blobUrl": AzureBlobStorage().get_blob_url(indexed_blob_name),
            }
            return document
        except Exception as e:
            raise Exception(f"Error embedding chunk {i}: {e}")

    # チャンクを5つのバッチに分割
    total_chunks = len(chunks)
    batch_size = min(batch_size, total_chunks)  # バッチサイズは最大でチャンク数まで
    chunk_batches = [[] for _ in range(batch_size)]

    # チャンクを均等に分配
    for i, chunk in enumerate(chunks):
        batch_index = i % batch_size
        chunk_batches[batch_index].append((i, chunk))

    # 各バッチを並列処理
    all_documents = []

    def process_batch(batch):
        batch_documents = []
        for i, chunk in batch:
            try:
                document = _embed_and_prepare_document(
                    open_ai_client, chunk, source_file_name, indexed_blob_name, i
                )
                batch_documents.append(document)
            except Exception as e:
                logging.error(f"Error processing chunk {i}: {e}")
        return batch_documents

    with ThreadPoolExecutor(max_workers=5) as executor:
        futures = [executor.submit(process_batch, batch) for batch in chunk_batches]
        for future in as_completed(futures):
            all_documents.extend(future.result())

    # ドキュメントをアップロード（最大100件ずつ）
    upload_batch_size = 100
    for i in range(0, len(all_documents), upload_batch_size):
        batch = all_documents[i : i + upload_batch_size]
        if batch:
            search_client.upload_documents(documents=batch)
            logging.info(f"Uploaded {len(batch)} documents to index {index_name}")


def index_pdf_docs(
    fileBytes: bytes, fileName: str, index_type_id: str, indexed_blob_name: str
):
    try:
        content = extract_markdown_text_from_pdf(fileBytes)
        chunks = _semantic_chunk(content)
        _index_docs_to_azure_ai_search(
            chunks, fileName, index_type_id, indexed_blob_name
        )
    except Exception as e:
        logging.error(f"Error indexing PDF: {e}\n{traceback.format_exc()}")


def index_docx_docs(
    fileBytes: bytes, fileName: str, index_type_id: str, indexed_blob_name: str
):
    try:
        content = extract_markdown_text_from_docx(fileBytes)
        chunks = _semantic_chunk(content)
        _index_docs_to_azure_ai_search(
            chunks, fileName, index_type_id, indexed_blob_name
        )
    except Exception as e:
        logging.error(f"Error indexing DOCX: {e}\n{traceback.format_exc()}")


def index_pptx_docs(
    fileBytes: bytes, fileName: str, index_type_id: str, indexed_blob_name: str
):
    try:
        content = extract_markdown_text_from_pptx(fileBytes)
        chunks = _semantic_chunk(content)
        _index_docs_to_azure_ai_search(
            chunks, fileName, index_type_id, indexed_blob_name
        )
    except Exception as e:
        logging.error(f"Error indexing PPTX: {e}\n{traceback.format_exc()}")


def index_html_docs(
    fileBytes: bytes, fileName: str, index_type_id: str, indexed_blob_name: str
):
    try:
        content = extract_markdown_text_from_html(fileBytes)
        chunks = _semantic_chunk(content)
        _index_docs_to_azure_ai_search(
            chunks, fileName, index_type_id, indexed_blob_name
        )
    except Exception as e:
        logging.error(f"Error indexing HTML: {e}\n{traceback.format_exc()}")


def index_image_docs(
    fileBytes: bytes, fileName: str, index_type_id: str, indexed_blob_name: str
):
    try:
        content = extract_markdown_text_from_image(fileBytes)
        chunks = _semantic_chunk(content)
        _index_docs_to_azure_ai_search(
            chunks, fileName, index_type_id, indexed_blob_name
        )
    except Exception as e:
        logging.error(f"Error indexing Image: {e}\n{traceback.format_exc()}")


def index_excel_docs(
    fileBytes: bytes, fileName: str, index_type_id: str, indexed_blob_name: str
):
    try:
        content = extract_markdown_text_from_excel(fileBytes)
        chunks = _semantic_chunk(content)
        _index_docs_to_azure_ai_search(
            chunks, fileName, index_type_id, indexed_blob_name
        )
    except Exception as e:
        logging.error(f"Error indexing Excel: {e}\n{traceback.format_exc()}")


def index_videostep_docs(
    fileBytes: bytes, fileName: str, index_type_id: str, indexed_blob_name: str
):
    try:
        import json

        content = json.loads(fileBytes)
        chunks = _semantic_chunk(content)
        _index_docs_to_azure_ai_search(
            chunks, fileName, index_type_id, indexed_blob_name
        )
    except Exception as e:
        logging.error(f"Error indexing videostep docs: {e}\n{traceback.format_exc()}")


def _semantic_chunk(contents: list) -> list:
    headers_to_split_on = [
        ("#", "Header 1"),
        ("##", "Header 2"),
        ("###", "Header 3"),
    ]
    from langchain.text_splitter import (
        MarkdownHeaderTextSplitter,
        RecursiveCharacterTextSplitter,
    )

    # まずMarkdownHeaderTextSplitterでヘッダーに基づいて分割
    text_splitter = MarkdownHeaderTextSplitter(headers_to_split_on=headers_to_split_on)
    content_text = "".join([c["page_content"] for c in contents])
    chunks = text_splitter.split_text(content_text)

    # 次に、各チャンクをさらに小さく分割（トークン数制限に合わせる）
    # 最大トークン数を8000に設定（安全マージンを取る）
    recursive_splitter = RecursiveCharacterTextSplitter(
        chunk_size=4000,  # 文字数で約4000文字（約4000トークンに相当）
        chunk_overlap=200,
        length_function=len,
    )

    chunks_with_page_number = []
    page_number = 1

    for chunk in chunks:
        keywords = ""
        if "Header 1" in chunk.metadata:
            keywords += chunk.metadata["Header 1"]
        if "Header 2" in chunk.metadata:
            keywords += " " + chunk.metadata["Header 2"]
        if "Header 3" in chunk.metadata:
            keywords += " " + chunk.metadata["Header 3"]

        # ページ番号の処理
        current_page_number = page_number
        matches = re.findall(r"<PAGE_NUMBER>(.*?)</PAGE_NUMBER>", chunk.page_content)
        if matches:
            page_number = int(matches[-1]) + 1
            chunk.page_content = re.sub(
                r"<PAGE_NUMBER>(.*?)</PAGE_NUMBER>", "", chunk.page_content
            )

        # 大きなチャンクをさらに小さく分割
        if len(chunk.page_content) > 4000:  # 文字数が4000を超える場合は分割
            smaller_chunks = recursive_splitter.split_text(chunk.page_content)
            for i, smaller_chunk in enumerate(smaller_chunks):
                chunks_with_page_number.append(
                    {
                        "content": smaller_chunk,
                        "page_number": current_page_number,
                        "keywords": keywords,
                    }
                )
        else:
            chunks_with_page_number.append(
                {
                    "content": chunk.page_content,
                    "page_number": current_page_number,
                    "keywords": keywords,
                }
            )

    return chunks_with_page_number


# ====================================================
# Blob Processing Functions
# ====================================================
def determine_file_type(filename):
    ext = os.path.splitext(filename)[1].lower()
    if ext == ".pdf":
        return "pdf"
    elif ext in [".doc", ".docx"]:
        return "word"
    elif ext in [".ppt", ".pptx"]:
        return "powerpoint"
    elif ext in [".html", ".htm"]:
        return "html"
    elif ext in [".xls", ".xlsx"]:
        return "excel"
    elif ext in [".jpg", ".jpeg", ".png", ".gif"]:
        return "image"
    else:
        return "unknown"


def extract_index_type(blob_name: str) -> str:
    """
    Blob 名から "SampleForIndex20250122/" 配下のフォルダ名を抽出して返す。
    例: "SampleForIndex20250122/FolderA/file.txt" -> "FolderA"
    """
    prefix = "SampleForIndex20250122/"
    if blob_name.startswith(prefix):
        parts = blob_name.split("/")
        if len(parts) >= 2:
            return parts[1]
    return None


def process_and_index_blob(blob_name, blob_service_client):
    source_container_client = blob_service_client.get_container_client(SOURCE_CONTAINER)
    blob_client = source_container_client.get_blob_client(blob_name)
    download_stream = blob_client.download_blob()
    file_bytes = download_stream.readall()

    ext = os.path.splitext(blob_name)[1].lower()
    file_type = determine_file_type(blob_name)
    index_type = extract_index_type(blob_name)
    if index_type is None:
        logging.error(
            f"Blob {blob_name} is not under the expected directory structure."
        )
        return None

    if file_type == "word" and ext == ".doc":
        try:
            file_bytes = convert_doc_bytes_to_docx_bytes(file_bytes)
            blob_name = os.path.splitext(blob_name)[0] + ".docx"
            logging.info(f"Converted .doc to .docx for {blob_name}")
        except Exception as e:
            logging.error(f"Error converting DOC to DOCX for {blob_name}: {e}")
            return None

    if file_type == "image":
        try:
            file_bytes = convert_image_to_pdf(file_bytes)
            file_type = "pdf"
            logging.info(f"Converted image to PDF for {blob_name}")
        except Exception as e:
            logging.error(f"Error converting image to PDF for {blob_name}: {e}")
            return None

    # フォルダ名からindex_type_idを取得
    db = SessionLocal()
    try:
        index_type_record = (
            db.query(SearchIndexType)
            .filter(SearchIndexType.folder_name == index_type)
            .first()
        )
        if not index_type_record:
            logging.error(f"No search index type found for folder: {index_type}")
            return None
        index_type_id = index_type_record.id
    finally:
        db.close()

    new_ext = os.path.splitext(blob_name)[1]
    new_blob_name = f"{uuid.uuid4()}{new_ext}"

    destination_container_client = blob_service_client.get_container_client(
        DESTINATION_CONTAINER
    )
    try:
        destination_container_client.create_container()
    except Exception:
        pass
    destination_blob_client = destination_container_client.get_blob_client(
        new_blob_name
    )
    destination_blob_client.upload_blob(file_bytes, overwrite=True)
    logging.info(f"File {blob_name} processed and uploaded as {new_blob_name}")

    if file_type == "pdf":
        index_pdf_docs(file_bytes, blob_name, index_type_id, new_blob_name)
    elif file_type == "word":
        index_docx_docs(file_bytes, blob_name, index_type_id, new_blob_name)
    elif file_type == "powerpoint":
        index_pptx_docs(file_bytes, blob_name, index_type_id, new_blob_name)
    elif file_type == "html":
        index_html_docs(file_bytes, blob_name, index_type_id, new_blob_name)
    elif file_type == "excel":
        index_excel_docs(file_bytes, blob_name, index_type_id, new_blob_name)
    elif file_type == "image":
        index_image_docs(file_bytes, blob_name, index_type_id, new_blob_name)
    elif file_type == "videostep":
        index_videostep_docs(file_bytes, blob_name, index_type_id, new_blob_name)
    else:
        logging.error(f"Unsupported file type for blob: {blob_name}")
        return None

    return file_type, index_type_id, new_blob_name


def process_new_blobs():
    blob_service_client = BlobServiceClient.from_connection_string(
        AZURE_STORAGE_CONNECTION_STRING
    )
    source_container_client = blob_service_client.get_container_client(SOURCE_CONTAINER)

    # まず、対象ディレクトリ ("SampleForIndex20250122/") 配下のフォルダ名を更新する
    # この処理は並列処理の前に一度だけ実行
    update_search_index_types(blob_service_client)

    # 処理対象のBlobリストを取得
    blobs_to_process = []
    db = SessionLocal()
    try:
        for blob in source_container_client.list_blobs():
            existing = (
                db.query(IndexedFile)
                .filter(IndexedFile.original_blob_name == blob.name)
                .first()
            )
            if not existing:
                blobs_to_process.append(blob.name)
                logging.info(f"Found new blob to process: {blob.name}")
    finally:
        db.close()

    if not blobs_to_process:
        logging.info("No new blobs to process.")
        return

    # 並列処理用の関数
    def process_blob(blob_name):
        try:
            logging.info(f"Processing blob: {blob_name}")
            result = process_and_index_blob(blob_name, blob_service_client)
            if result is None:
                logging.warning(f"Failed to process blob: {blob_name}")
            else:
                logging.info(f"Successfully processed blob: {blob_name}")
            return blob_name, result
        except Exception as e:
            logging.error(
                f"Error processing blob {blob_name}: {e}\n{traceback.format_exc()}"
            )
            return blob_name, None

    # 最大5スレッドで並列処理
    results = []
    with ThreadPoolExecutor(max_workers=5) as executor:
        futures = {
            executor.submit(process_blob, blob_name): blob_name
            for blob_name in blobs_to_process
        }
        for future in as_completed(futures):
            blob_name = futures[future]
            try:
                _, result = future.result()
                if result is not None:
                    results.append((blob_name, result))
            except Exception as e:
                logging.error(
                    f"Exception processing blob {blob_name}: {e}\n{traceback.format_exc()}"
                )

    logging.info(f"Processed {len(results)} out of {len(blobs_to_process)} blobs")


# ====================================================
# Update Search Index Types & Create Index
# ====================================================
from sqlalchemy.exc import SQLAlchemyError
from azure.search.documents.indexes.models import (
    SimpleField,
    SearchableField,
    SearchField,
    SearchFieldDataType,
    SearchIndex,
    SemanticSearch,
    SemanticConfiguration,
    SemanticPrioritizedFields,
    SemanticField,
    VectorSearch,
    HnswAlgorithmConfiguration,
    VectorSearchProfile,
)


def create_index(index_name: str):
    client = AzureAISearch().init_search_index_client()
    if index_name in client.list_index_names():
        logging.info(f"Index {index_name} already exists.")
        return

    fields = [
        SimpleField(name="id", type=SearchFieldDataType.String, key=True),
        SimpleField(name="keywords", type="Edm.String", analyzer_name="ja.microsoft"),
        SearchableField(
            name="content", type="Edm.String", analyzer_name="ja.microsoft"
        ),
        SearchField(
            name="contentVector",
            type=SearchFieldDataType.Collection(SearchFieldDataType.Single),
            searchable=True,
            vector_search_dimensions=1536,
            vector_search_profile_name="myHnswProfile",
        ),
        SimpleField(name="sourceFileName", type=SearchFieldDataType.String),
        SimpleField(name="pageNumber", type=SearchFieldDataType.Int32),
        SimpleField(name="blobUrl", type=SearchFieldDataType.String),
    ]
    semantic_settings = SemanticSearch(
        configurations=[
            SemanticConfiguration(
                name="default",
                prioritized_fields=SemanticPrioritizedFields(
                    title_field=None,
                    content_fields=[SemanticField(field_name="content")],
                ),
            )
        ]
    )
    vector_search = VectorSearch(
        algorithms=[HnswAlgorithmConfiguration(name="myHnsw")],
        profiles=[
            VectorSearchProfile(
                name="myHnswProfile",
                algorithm_configuration_name="myHnsw",
            )
        ],
    )
    index = SearchIndex(
        name=index_name,
        fields=fields,
        vector_search=vector_search,
        semantic_search=semantic_settings,
    )
    client.create_index(index)
    logging.info(f"Created search index: {index_name}")


def update_search_index_types(blob_service_client):
    prefix = "SampleForIndex20250122/"
    container_client = blob_service_client.get_container_client(SOURCE_CONTAINER)
    folder_set = set()

    # フォルダ名を収集
    for blob in container_client.list_blobs(name_starts_with=prefix):
        parts = blob.name.split("/")
        if len(parts) >= 2:
            folder_set.add(parts[1])

    if not folder_set:
        logging.info("No folders found under the specified prefix.")
        return

    # データベースから既存のフォルダ名を取得
    session = SessionLocal()
    try:
        existing_folders = {
            record.folder_name for record in session.query(SearchIndexType).all()
        }

        # 新しいフォルダのみを処理
        new_folders = folder_set - existing_folders
        if not new_folders:
            logging.info("No new folders to add.")
            return

        logging.info(f"Found {len(new_folders)} new folders to add: {new_folders}")

        # 新しいフォルダを追加
        for folder in new_folders:
            try:
                # トランザクションを使用して重複を防止
                with session.begin_nested():
                    # 再度確認（他のプロセスが同時に挿入した可能性）
                    exists = (
                        session.query(SearchIndexType)
                        .filter(SearchIndexType.folder_name == folder)
                        .first()
                    )
                    if not exists:
                        # 新しいUUIDを生成
                        index_id = str(uuid.uuid4())
                        # 新しいレコードを作成し、生成したUUIDをidとして設定
                        new_type = SearchIndexType(id=index_id, folder_name=folder)
                        session.add(new_type)
                        # コミットして同じUUIDをインデックス名として使用
                        session.flush()
                        create_index(index_id)
                        logging.info(
                            f"Inserted new search index type: {folder} with index id: {index_id}"
                        )
            except SQLAlchemyError as e:
                logging.error(
                    f"Error inserting search index type for folder {folder}: {e}"
                )
                # このフォルダの処理をスキップして次へ
                continue

        # 最終的にコミット
        session.commit()
    except SQLAlchemyError as e:
        session.rollback()
        logging.error(f"Error updating search index types: {e}")
    finally:
        session.close()


# ====================================================
# 新規 Blob ファイル一覧の取得（DB と比較）＋ tqdm で進捗表示
# ====================================================
def list_new_blob_files():
    blob_service_client = BlobServiceClient.from_connection_string(
        AZURE_STORAGE_CONNECTION_STRING
    )
    container_client = blob_service_client.get_container_client(SOURCE_CONTAINER)
    blob_names = [blob.name for blob in container_client.list_blobs()]

    db = SessionLocal()
    try:
        existing_files = {
            record.original_blob_name for record in db.query(IndexedFile).all()
        }
    finally:
        db.close()

    new_files = [name for name in blob_names if name not in existing_files]

    for file in tqdm(new_files, desc="新規Blobファイル一覧"):
        logging.info(f"New file found: {file}")

    return new_files


# ====================================================
# Azure Functions エントリポイント
# ====================================================
app = azfunc.FunctionApp()


@app.function_name(name="BlobIndexerFunction")
@app.schedule(schedule="0 0 0 * * *", arg_name="mytimer", run_on_startup=True)
def BlobIndexerFunction(mytimer: azfunc.TimerRequest) -> None:
    logging.info("Azure Function Triggered for Blob Indexing")
    new_files = list_new_blob_files()
    if new_files:
        logging.info("新規追加されたファイル一覧:")
        for file in new_files:
            logging.info(file)
    process_new_blobs()

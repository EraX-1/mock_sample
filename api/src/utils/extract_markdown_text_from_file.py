import io
import os
import tempfile
from itertools import groupby

import html2text
from bs4 import BeautifulSoup
from docx import Document
from openpyxl import Workbook, load_workbook
from pptx import Presentation
from pypdf import PdfReader, PdfWriter
from tqdm import tqdm

from src.services.azure_ai_doc_intel import AzureAIDocumentIntelligence
from src.utils.convert_file_to_pdf import convert_image_to_pdf


def extract_markdown_text_from_docx(file: bytes) -> list[dict[str, any]]:
    """Wordファイル(.doc, .docx)のbytes型データからマークダウン形式でテキストを抽出する

    Args:
        file (bytes): Wordファイル(.doc, .docx)のbytes型データ

    Returns:
        list[dict[str, any]]: 各ページのテキストとメタデータを含む辞書のリスト
    """
    paragraph_contents = []

    def process_paragraph(i_para):
        """各段落を処理する"""
        i, paragraph = i_para
        try:
            # TODO: 1ページの基準を定義してページ番号を付与する必要がある
            # 段落のテキストを取得
            text = paragraph.text.strip()
            if text:
                markdown_text = ""

                # 段落スタイルに基づいてマークダウン形式に変換
                if paragraph.style.name.startswith("Heading"):
                    level = int(paragraph.style.name[-1])  # Heading 1, 2, 3などの数字を取得
                    markdown_text = "#" * level + " " + text + "\n\n"  # ヘッダー後に2つの改行を追加
                else:
                    markdown_text = text + "\n\n"  # 通常の段落も2つの改行で区切る

                return {
                    "content": markdown_text,
                    "page": 1,  # Wordファイルではページ番号の取得が難しいため固定値
                    "order": i,  # 段落の順序を保持
                }
            return None
        except Exception as e:
            raise Exception(f"段落の処理中にエラーが発生しました {i}: {e}")

    # Wordファイルを読み込む
    byte_stream = io.BytesIO(file)
    doc = Document(byte_stream)

    # NOTE: 非同期処理ではパフォーマンスが上がらず、並列・並行処理は上手く実装できなかった
    for i, para in enumerate(doc.paragraphs):
        res = process_paragraph((i, para))
        if res is not None:
            paragraph_contents.append(res)

    # 段落の順序でソートする
    paragraph_contents.sort(key=lambda x: x["order"])

    # ページごとにコンテンツをまとめる
    docs = []
    for page, group in groupby(paragraph_contents, key=lambda x: x["page"]):
        combined_content = ""
        for item in group:
            combined_content += item["content"]

        docs.append(
            {
                "page_content": combined_content + f"<PAGE_NUMBER>{page}</PAGE_NUMBER>",
                "metadata": {"page": page},
            }
        )

    return docs


def extract_markdown_text_from_pptx(file: bytes) -> list[dict[str, any]]:
    """PowerPointファイル(.pptx)のbytes型データからマークダウン形式でテキストを抽出する

    Args:
        file (bytes): PowerPointファイル(.pptx)のbytes型データ

    Returns:
        list[dict[str, any]]: 各ページのテキストとメタデータを含む辞書のリスト
    """
    docs = []

    def process_slide(i_slide):
        """各スライドを処理する"""
        i, slide = i_slide
        slide_docs = []
        try:
            # スライドのテキストを抽出
            texts = []

            # タイトルの抽出
            if slide.shapes.title:
                texts.append(f"# {slide.shapes.title.text}\n")

            # その他のテキストの抽出
            for shape in slide.shapes:
                if (
                    hasattr(shape, "text")
                    and shape.text.strip()
                    and shape != slide.shapes.title
                ):
                    # テキストボックスの場合
                    text = shape.text.strip()
                    if text:
                        texts.append(text + "\n")

            # スライドから抽出したテキストを結合
            content = "\n".join(texts).strip()

            if content:
                slide_docs.append(
                    {"page_content": content, "metadata": {"page": i + 1}}
                )

            return slide_docs

        except Exception as e:
            raise Exception(f"スライドの処理中にエラーが発生しました {i}: {e}")

    temp_file = tempfile.NamedTemporaryFile(delete=False)
    try:
        temp_file.write(file)
        temp_file.close()

        # PowerPointファイルを読み込む
        prs = Presentation(temp_file.name)

        # NOTE: 非同期処理ではパフォーマンスが上がらず、並列・並行処理は上手く実装できなかった
        for i, slide in enumerate(prs.slides):
            docs.extend(process_slide((i, slide)))

    finally:
        os.unlink(temp_file.name)

    # ページ番号順にソートする
    docs = sorted(docs, key=lambda x: x["metadata"]["page"])

    # ページ番号をテキストに埋め込む
    for doc in docs:
        doc["page_content"] += (
            "<PAGE_NUMBER>" + str(doc["metadata"]["page"]) + "</PAGE_NUMBER>"
        )

    return docs


def extract_markdown_text_from_html(file: bytes) -> list[dict[str, any]]:
    """HTMLファイル(.html)のbytes型データからマークダウン形式でテキストを抽出する

    Args:
        file (bytes): HTMLファイル(.html)のbytes型データ

    Returns:
        list[dict[str, any]]: 各ページのテキストとメタデータを含む辞書のリスト
    """
    docs = []

    def process_html_section(i_section):
        """各セクションを処理する"""
        i, section = i_section
        section_docs = []
        try:
            h2t = html2text.HTML2Text()
            h2t.ignore_links = False
            h2t.ignore_images = False
            h2t.ignore_tables = False

            # セクションのテキストを抽出してマークダウンに変換
            markdown_content = h2t.handle(str(section))

            # 空のセクションをスキップ
            if markdown_content.strip():
                section_docs.append(
                    {"page_content": markdown_content, "metadata": {"page": i + 1}}
                )

            return section_docs

        except Exception as e:
            raise Exception(f"セクションの処理中にエラーが発生しました {i}: {e}")

    try:
        # HTMLコンテンツをパース
        soup = BeautifulSoup(file.decode("utf-8"), "html.parser")

        # コンテンツを意味のある単位で分割
        # article, section, div などの主要なコンテナ要素を検索
        sections = soup.find_all(
            ["article", "section", "div"], class_=True, recursive=False
        )
        if not sections:
            # セクションが見つからない場合は、bodyを1つのセクションとして扱う
            sections = [soup.body] if soup.body else [soup]

        # NOTE: 非同期処理ではパフォーマンスが上がらず、並列・並行処理は上手く実装できなかった
        for i, section in enumerate(sections):
            docs.extend(process_html_section((i, section)))

    except Exception as e:
        raise Exception(f"Error processing HTML file: {e}")

    # ページ番号順にソートする
    docs = sorted(docs, key=lambda x: x["metadata"]["page"])

    # ページ番号をテキストに埋め込む
    for doc in docs:
        doc["page_content"] += f"<PAGE_NUMBER>{doc['metadata']['page']}</PAGE_NUMBER>"

    return docs


def extract_markdown_text_from_pdf(file: bytes) -> list[dict[str, any]]:
    """PDFファイル(.pdf)のbytes型データからマークダウン形式でテキストを抽出する

    Args:
        file (bytes): PDFファイル(.pdf)のbytes型データ

    Returns:
        list[dict[str, any]]: 各ページのテキストとメタデータを含む辞書のリスト
    """
    docs = []

    def process_page(i_page):
        i, page = i_page
        page_docs = []
        temp_page_file = None
        try:
            temp_page_file = tempfile.NamedTemporaryFile(delete=False)
            writer = PdfWriter()
            writer.add_page(page)
            # メモリ上のバッファにPDFを書き込む
            bytes_io = io.BytesIO()
            writer.write(bytes_io)
            # バイト列を取得
            page_bytes = bytes_io.getvalue()
            bytes_io.close()  # リソースを解放
            # バイト列をファイルに書き込む
            temp_page_file.write(page_bytes)
            temp_page_file.close()  # 書き込み終了後に閉じる

            loader = AzureAIDocumentIntelligence().init_loader(
                file_path=temp_page_file.name
            )
            doc = loader.load()
            for d in doc:
                # Document Intelligence のエラーレスポンスを検出
                if (
                    "サポートされていないファイル形式です" in d.page_content
                    or d.page_content.strip() == ""
                ):
                    raise Exception(
                        "Document Intelligence returned error or empty content"
                    )

                page_docs.append(
                    {
                        "page_content": d.page_content,
                        "metadata": {"page": i + 1},
                    }
                )

            return page_docs

        except Exception as e:
            print(f"⚠️ Document Intelligence failed for page {i}: {e}")
            print("🔄 Falling back to simple PDF text extraction")
            # フォールバック: シンプルなPDF text extraction
            try:
                text = page.extract_text()
                if text.strip():
                    page_docs.append(
                        {
                            "page_content": f"# Page {i+1}\n\n{text}",
                            "metadata": {"page": i + 1},
                        }
                    )
                else:
                    page_docs.append(
                        {
                            "page_content": f"# Page {i+1}\n\n(No text content found)",
                            "metadata": {"page": i + 1},
                        }
                    )
                return page_docs
            except Exception as fallback_e:
                print(f"❌ Fallback also failed: {fallback_e}")
                # 最後の手段として空のコンテンツを返す
                page_docs.append(
                    {
                        "page_content": f"# Page {i+1}\n\n(Text extraction failed)",
                        "metadata": {"page": i + 1},
                    }
                )
                return page_docs

        finally:
            if temp_page_file is not None:
                os.unlink(temp_page_file.name)

    temp_file = tempfile.NamedTemporaryFile(delete=False)
    try:
        temp_file.write(file)
        temp_file.close()  # 書き込み終了後に閉じ
        reader = PdfReader(temp_file.name)

        # NOTE: 非同期処理ではパフォーマンスが上がらず、並列・並行処理は上手く実装できなかった
        for i, page in enumerate(reader.pages):
            docs.extend(process_page((i, page)))

    finally:
        os.unlink(temp_file.name)

    # ページ番号順にソートする
    docs = sorted(docs, key=lambda x: x["metadata"]["page"])

    # ページ番号をテキストに埋め込む
    for doc in docs:
        doc["page_content"] += (
            "<PAGE_NUMBER>" + str(doc["metadata"]["page"]) + "</PAGE_NUMBER>"
        )

    return docs


def extract_markdown_text_from_image(file: bytes) -> list[dict[str, any]]:
    """画像データ(.png, .jpg, .jpeg)のbytes型データからマークダウン形式でテキストを抽出する

    Args:
        file (bytes): 画像(.png, .jpg, .jpeg)のbytesデータ

    Returns:
        list[dict[str, any]]: 各ページのテキストとメタデータを含む辞書のリスト
    """
    # 画像データからPDFデータへ変換
    pdf_bytes = convert_image_to_pdf(file)

    # PDFデータからテキスト抽出
    docs = extract_markdown_text_from_pdf(pdf_bytes)

    return docs


def extract_markdown_text_from_excel(file: bytes):
    """Excelファイル(.xlsx, .xls, .xlsm)のbytes型データからマークダウン形式でテキストを抽出する

    Args:
        file (bytes): Excelファイル(.xlsx, .xls, .xlsm)のbytes型データ

    Returns:
        list[dict[str, any]]: 各ページのテキストとメタデータを含む辞書のリスト
    """
    # 既存のExcelファイルを読み込み
    # BytesIOオブジェクトに変換して読み込む
    excel_file = io.BytesIO(file)
    # ワークブック読み込み時にマクロを読み込まない
    wb = load_workbook(excel_file, keep_vba=False)
    # VBA関連の属性をクリア
    wb.vba_archive = None
    wb.is_template = False
    wb.vba_modified = False

    # シートごとのReadableBuffer（BytesIOオブジェクト）を格納するリスト
    docs = []

    for i, sheet in enumerate(wb.worksheets):
        print(f"シート名: {sheet.title}")
        # 新しいワークブックを作成し、現在のシートの内容をコピー
        new_wb = Workbook()
        new_ws = new_wb.active
        new_ws.title = sheet.title

        none_count = 0
        max_rows = sheet.max_row
        max_cols = sheet.max_column

        for row_idx in tqdm(range(1, max_rows + 1)):
            none_row = True
            for col_idx in range(1, max_cols + 1):
                cell = sheet.cell(row=row_idx, column=col_idx)
                value = cell.value
                if value is not None:
                    none_row = False
                new_ws.cell(row=row_idx, column=col_idx, value=value)
            if none_row:
                none_count += 1
            else:
                none_count = 0
            # 10行連続でデータが無かったらbreak
            if none_count > 10:
                break

        # シートごとにBytesIOに保存
        sheet_buffer = io.BytesIO()
        new_wb.save(sheet_buffer)
        sheet_buffer.seek(0)

        with tempfile.NamedTemporaryFile(delete=False, suffix=".xlsx") as temp_file:
            temp_file.write(sheet_buffer.getvalue())
            temp_file.flush()
            loader = AzureAIDocumentIntelligence().init_loader(file_path=temp_file.name)
            doc = loader.load()

            docs.append(
                {
                    "page_content": doc[0].page_content,
                    "metadata": {"page": i + 1},
                }
            )

    return docs
